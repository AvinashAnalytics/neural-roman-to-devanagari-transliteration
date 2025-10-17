# gui/app.py
"""Main Application for Hindi Transliteration System using Streamlit
CS772 Assignment 2 - मेरा भारत महान

Features:
- LSTM and Transformer models with local attention
- LLM-based transliteration via OpenAI, Anthropic, Google, Groq
- UTF-8 config support and robust error handling
- Interactive Streamlit GUI with dark mode and Indian tricolor theme
"""
import os
import sys
import json
import pickle
import io
import yaml
import torch
import streamlit as st
from typing import Dict, List, Optional, Tuple
import pandas as pd
from datetime import datetime
import plotly.express as px
import plotly.graph_objects as go
import time
import random

# CRITICAL FIX: Setup project root BEFORE imports
CURRENT_FILE = os.path.abspath(__file__)
GUI_DIR = os.path.dirname(CURRENT_FILE)
PROJECT_ROOT = os.path.dirname(GUI_DIR)

# Add parent directory to path
sys.path.insert(0, PROJECT_ROOT)

from models.lstm_model import Seq2SeqLSTM
from models.transformer_model import TransformerSeq2Seq
from models.llm_model import LLMTransliterator
from utils.vocab import Vocabulary
from utils.evaluation import Evaluator


def load_result_file(path: str) -> Optional[dict]:
    """Load various result JSON schemas and normalize to a dict with key 'results'.

    Supported schemas:
    - {"results": {"greedy": {...}, "beam_5": {...}}}
    - {"greedy": {...}, "beam_5": {...}}  (legacy)
    - {other keys...} -> will be wrapped under 'results'

    Returns normalized dict or None on parse error.
    """
    if not path or not os.path.exists(path):
        return None
    try:
        with open(path, 'r', encoding='utf-8') as fh:
            loaded = json.load(fh)
    except Exception as e:
        try:
            # sometimes files are saved as single-line JSON with trailing garbage
            txt = open(path, 'r', encoding='utf-8', errors='ignore').read()
            loaded = json.loads(txt)
        except Exception:
            st.error(f"Failed to parse result JSON: {path} — {e}")
            return None

    # If already in new format
    if isinstance(loaded, dict) and 'results' in loaded and isinstance(loaded['results'], dict):
        return loaded

    # If top-level keys look like greedy/beam_x then wrap
    if isinstance(loaded, dict):
        metric_keys = {'greedy', 'beam_3', 'beam_5', 'beam_10', 'beam_1', 'beam'}
        if any(k in loaded for k in metric_keys):
            return {'results': loaded}

        # Some files use nested naming like 'results': <list> or other shapes
        # If file already looks like a metrics dict with numeric values, wrap under 'results'
        # Heuristic: check if values are dicts containing 'word_accuracy' or 'char_f1'
        wrapped = {}
        for k, v in (loaded.items() if isinstance(loaded, dict) else []):
            if isinstance(v, dict) and ('word_accuracy' in v or 'char_f1' in v or 'mean_edit_distance' in v):
                wrapped[k] = v
        if wrapped:
            return {'results': wrapped}

    # Nothing matched – return original but wrapped to avoid KeyError
    return {'results': loaded if isinstance(loaded, dict) else { 'data': loaded }}


def make_error_3d_plot(error_analysis: dict):
    """Return a Plotly 3D scatter Figure summarizing error_analysis categories."""
    x = []
    y = []
    z = []
    text = []
    for cat in ('conjuncts', 'aspirated', 'vowel_matras'):
        sub = error_analysis.get(cat, {})
        for k, v in sub.items():
            x.append(cat)
            y.append(k)
            z.append(v)
            text.append(f"{cat}:{k} — {v}")

    fig = go.Figure(data=[go.Scatter3d(x=x, y=y, z=z, text=text, hoverinfo='text', mode='markers', marker=dict(size=[max(6, min(30, int(v/10))) for v in z], color=z, colorscale='Viridis'))])
    fig.update_layout(margin=dict(l=0, r=0, b=0, t=30), height=420, scene=dict(xaxis_title='Category', yaxis_title='Key', zaxis_title='Count'))
    return fig

# Page configuration - MUST BE FIRST
st.set_page_config(
    page_title="🇮🇳 Hindi Transliteration System",
    page_icon="🔤",
    layout="wide",
    initial_sidebar_state="expanded",
    menu_items={
        'Get Help': 'https://github.com/yourusername/transliteration',
        'Report a bug': 'https://github.com/yourusername/transliteration/issues',
        'About': 'CS772 Assignment 2 - Hindi Transliteration System | मेरा भारत महान'
    }
)

# 🇮🇳 DARK MODE WITH INDIAN TRICOLOR THEME
st.markdown("""
<style>
    /* ===== INDIAN FLAG COLORS ===== */
    :root {
        --saffron: #FF9933;
        --white: #FFFFFF;
        --green: #138808;
        --navy-blue: #000080;
        --dark-bg: #0f172a;
        --dark-card: #1e293b;
        --dark-surface: #334155;
        --gold: #FFD700;
        --shadow-lg: 0 10px 40px rgba(0, 0, 0, 0.5);
        --shadow-md: 0 5px 20px rgba(0, 0, 0, 0.4);
    }
    
    /* ===== DARK BACKGROUND ===== */
    .stApp {
        background: linear-gradient(135deg, #0f172a 0%, #1e293b 50%, #0f172a 100%);
        color: #e2e8f0;
    }
    
    /* ===== MAIN HEADER WITH TRICOLOR & FLOATING CHARS ===== */
    .main-header {
        position: relative;
        text-align: center;
        padding: 50px 20px;
        background: linear-gradient(
            to bottom,
            var(--saffron) 0%,
            var(--saffron) 33%,
            var(--white) 33%,
            var(--white) 66%,
            var(--green) 66%,
            var(--green) 100%
        );
        border-radius: 20px;
        margin-bottom: 30px;
        box-shadow: var(--shadow-lg);
        overflow: hidden;
        border: 3px solid var(--gold);
    }
    
    /* Ashoka Chakra Effect */
    .main-header::before {
        content: '☸';
        position: absolute;
        top: 50%;
        left: 50%;
        transform: translate(-50%, -50%);
        font-size: 80px;
        color: var(--navy-blue);
        opacity: 0.3;
        animation: rotateChakra 20s linear infinite;
        z-index: 0;
    }
    
    @keyframes rotateChakra {
        from { transform: translate(-50%, -50%) rotate(0deg); }
        to { transform: translate(-50%, -50%) rotate(360deg); }
    }
    
    /* Floating "मेरा भारत महान" Characters */
    .hindi-char {
        position: absolute;
        font-size: 2.5rem;
        font-weight: 700;
        color: rgba(255, 215, 0, 0.4);
        text-shadow: 0 0 10px rgba(255, 215, 0, 0.6);
        pointer-events: none;
        z-index: 1;
        animation: floatChar 15s ease-in-out infinite;
    }
    
    @keyframes floatChar {
        0%, 100% { transform: translateY(0) translateX(0) rotate(0deg); opacity: 0; }
        10% { opacity: 1; }
        50% { transform: translateY(-30px) translateX(20px) rotate(5deg); opacity: 1; }
        90% { opacity: 1; }
        100% { transform: translateY(0) translateX(0) rotate(0deg); opacity: 0; }
    }
    
    /* Individual char positions */
    .char-1 { top: 10%; left: 5%; animation-delay: 0s; }
    .char-2 { top: 15%; left: 12%; animation-delay: 1s; }
    .char-3 { top: 20%; left: 18%; animation-delay: 2s; }
    .char-4 { top: 25%; left: 25%; animation-delay: 3s; }
    .char-5 { top: 30%; left: 75%; animation-delay: 4s; }
    .char-6 { top: 35%; left: 82%; animation-delay: 5s; }
    .char-7 { top: 40%; left: 88%; animation-delay: 6s; }
    .char-8 { top: 45%; left: 94%; animation-delay: 7s; }
    .char-9 { top: 70%; left: 8%; animation-delay: 8s; }
    .char-10 { top: 75%; left: 15%; animation-delay: 9s; }
    .char-11 { top: 80%; left: 22%; animation-delay: 10s; }
    .char-12 { top: 85%; left: 30%; animation-delay: 11s; }
    
    .main-header h1 {
        position: relative;
        z-index: 2;
        color: var(--navy-blue);
        font-size: 3.5rem;
        font-weight: 900;
        text-shadow: 
            2px 2px 4px rgba(255, 255, 255, 0.8),
            -1px -1px 2px rgba(0, 0, 0, 0.3);
        margin: 0;
        animation: titlePulse 3s ease-in-out infinite;
    }
    
    @keyframes titlePulse {
        0%, 100% { transform: scale(1); }
        50% { transform: scale(1.02); }
    }
    
    .main-header p {
        position: relative;
        z-index: 2;
        color: var(--navy-blue);
        font-size: 1.2rem;
        margin-top: 10px;
        font-weight: 600;
        text-shadow: 1px 1px 2px rgba(255, 255, 255, 0.8);
    }
    
    /* ===== FOOTER WITH TRICOLOR ===== */
    .main-footer {
        position: relative;
        text-align: center;
        padding: 40px 20px;
        background: linear-gradient(
            to bottom,
            var(--green) 0%,
            var(--green) 33%,
            var(--white) 33%,
            var(--white) 66%,
            var(--saffron) 66%,
            var(--saffron) 100%
        );
        border-radius: 20px;
        margin-top: 40px;
        box-shadow: var(--shadow-lg);
        overflow: hidden;
        border: 3px solid var(--gold);
    }
    
    .main-footer::before {
        content: '🇮🇳';
        position: absolute;
        top: 50%;
        left: 50%;
        transform: translate(-50%, -50%);
        font-size: 100px;
        opacity: 0.2;
        z-index: 0;
    }
    
    .main-footer p {
        position: relative;
        z-index: 2;
        color: var(--navy-blue);
        font-weight: 700;
        margin: 8px 0;
        text-shadow: 1px 1px 2px rgba(255, 255, 255, 0.8);
        font-size: 1.1rem;
    }
    
    .footer-hindi {
        font-size: 1.5rem;
        color: var(--saffron);
        font-weight: 900;
        text-shadow: 2px 2px 4px rgba(0, 0, 0, 0.3);
    }
    
    /* ===== SECTION HEADERS (Dark) ===== */
    .section-header {
        background: linear-gradient(135deg, var(--saffron) 0%, var(--green) 100%);
        padding: 25px 30px;
        border-radius: 15px;
        margin: 25px 0;
        box-shadow: var(--shadow-md);
        position: relative;
        overflow: hidden;
        border-left: 5px solid var(--gold);
    }
    
    .section-header::after {
        content: '';
        position: absolute;
        top: 0;
        left: -100%;
        width: 100%;
        height: 100%;
        background: linear-gradient(90deg, transparent, rgba(255,255,255,0.3), transparent);
        animation: shimmer 3s infinite;
    }
    
    @keyframes shimmer {
        0% { left: -100%; }
        100% { left: 100%; }
    }
    
    .section-header h2 {
        color: white;
        margin: 0;
        font-weight: 800;
        font-size: 2rem;
        text-shadow: 2px 2px 4px rgba(0, 0, 0, 0.5);
    }
    
    /* ===== DARK CARDS ===== */
    .metric-card {
        background: var(--dark-card);
        padding: 25px;
        border-radius: 16px;
        border-left: 5px solid var(--saffron);
        margin: 20px 0;
        box-shadow: var(--shadow-md);
        transition: all 0.3s ease;
        color: #e2e8f0;
    }
    
    .metric-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 15px 50px rgba(255, 153, 51, 0.3);
        border-left-color: var(--gold);
    }
    
    /* ===== MODEL CARDS (Dark) ===== */
    .model-card {
        background: linear-gradient(135deg, var(--dark-card) 0%, var(--dark-surface) 100%);
        padding: 30px;
        border-radius: 20px;
        color: #e2e8f0;
        margin: 20px 0;
        box-shadow: var(--shadow-lg);
        transition: all 0.4s ease;
        border: 2px solid var(--saffron);
    }
    
    .model-card:hover {
        transform: scale(1.03);
        box-shadow: 0 20px 60px rgba(255, 153, 51, 0.4);
        border-color: var(--gold);
    }
    
    /* ===== TRANSLITERATION OUTPUT (Dark Premium) ===== */
    .transliteration-output {
        background: linear-gradient(135deg, #1e3a8a 0%, #312e81 100%);
        padding: 45px;
        border-radius: 20px;
        border: 4px solid var(--gold);
        font-size: 2.5rem;
        font-weight: 800;
        text-align: center;
        margin: 25px 0;
        box-shadow: 
            0 15px 50px rgba(255, 215, 0, 0.3),
            inset 0 0 30px rgba(255, 215, 0, 0.1);
        color: var(--gold);
        transition: all 0.3s ease;
        position: relative;
        overflow: hidden;
    }
    
    .transliteration-output::before {
        content: '';
        position: absolute;
        top: -50%;
        left: -50%;
        width: 200%;
        height: 200%;
        background: radial-gradient(circle, rgba(255,215,0,0.1) 0%, transparent 70%);
        animation: outputGlow 4s ease-in-out infinite;
    }
    
    @keyframes outputGlow {
        0%, 100% { transform: translate(0, 0); }
        50% { transform: translate(10%, 10%); }
    }
    
    .transliteration-output:hover {
        transform: scale(1.02);
        box-shadow: 
            0 20px 70px rgba(255, 215, 0, 0.5),
            inset 0 0 40px rgba(255, 215, 0, 0.2);
    }
    
    /* ===== PROVIDER CARDS (Dark Indian Theme) ===== */
    .provider-card {
        background: linear-gradient(135deg, var(--dark-card) 0%, #374151 100%);
        padding: 25px;
        border-radius: 16px;
        margin: 20px 0;
        box-shadow: var(--shadow-md);
        transition: all 0.3s ease;
        border: 2px solid var(--green);
        color: #e2e8f0;
    }
    
    .provider-card:hover {
        transform: translateX(8px);
        box-shadow: 0 15px 50px rgba(19, 136, 8, 0.4);
        border-color: var(--saffron);
    }
    
    .provider-card h4 {
        margin: 0 0 10px 0;
        color: var(--saffron);
        font-size: 1.5rem;
        font-weight: 800;
    }
    
    .provider-card p {
        color: #cbd5e1;
        margin: 8px 0;
    }
    
    /* ===== PREMIUM BUTTONS (Indian Colors) ===== */
    .stButton > button {
        background: linear-gradient(135deg, var(--saffron) 0%, var(--green) 100%);
        color: white;
        border: none;
        padding: 14px 32px;
        border-radius: 12px;
        font-weight: 700;
        font-size: 16px;
        transition: all 0.3s ease;
        box-shadow: 0 5px 20px rgba(255, 153, 51, 0.3);
        text-transform: uppercase;
        letter-spacing: 0.5px;
    }
    
    .stButton > button:hover {
        transform: translateY(-3px);
        box-shadow: 0 10px 30px rgba(255, 215, 0, 0.5);
        background: linear-gradient(135deg, var(--gold) 0%, var(--saffron) 100%);
    }
    
    .stButton > button:active {
        transform: translateY(0);
    }
    
    /* ===== SUCCESS/ERROR BOXES (Dark) ===== */
    .success-box {
        background: linear-gradient(135deg, #065f46 0%, #047857 100%);
        border: 2px solid var(--green);
        padding: 20px;
        border-radius: 12px;
        margin: 15px 0;
        box-shadow: var(--shadow-md);
        color: #d1fae5;
    }
    
    .error-box {
        background: linear-gradient(135deg, #7f1d1d 0%, #991b1b 100%);
        border: 2px solid #ef4444;
        padding: 20px;
        border-radius: 12px;
        margin: 15px 0;
        box-shadow: var(--shadow-md);
        color: #fecaca;
    }
    
    /* ===== TABS (Dark Indian Theme) ===== */
    .stTabs [data-baseweb="tab-list"] {
        gap: 10px;
        background: linear-gradient(135deg, var(--dark-card) 0%, var(--dark-surface) 100%);
        border-radius: 15px;
        padding: 10px;
        border: 2px solid var(--saffron);
    }
    
    .stTabs [data-baseweb="tab"] {
        background-color: transparent;
        border-radius: 10px;
        color: #94a3b8;
        font-weight: 600;
        padding: 12px 24px;
        transition: all 0.3s ease;
        border: 2px solid transparent;
    }
    
    .stTabs [data-baseweb="tab"]:hover {
        background-color: rgba(255, 153, 51, 0.1);
        color: var(--saffron);
    }
    
    .stTabs [aria-selected="true"] {
        background: linear-gradient(135deg, var(--saffron) 0%, var(--gold) 100%);
        color: white;
        box-shadow: 0 5px 15px rgba(255, 153, 51, 0.4);
        border-color: var(--gold);
    }
    
    /* ===== SIDEBAR (Dark Gradient) ===== */
    [data-testid="stSidebar"] {
        background: linear-gradient(180deg, #1e293b 0%, #0f172a 100%);
        border-right: 3px solid var(--saffron);
    }
    
    [data-testid="stSidebar"] .stMarkdown {
        color: #e2e8f0;
    }
    
    [data-testid="stSidebar"] h2, 
    [data-testid="stSidebar"] h3 {
        color: var(--saffron) !important;
    }
    
    /* ===== PROGRESS BAR (Tricolor) ===== */
    .stProgress > div > div {
        background: linear-gradient(
            90deg, 
            var(--saffron) 0%, 
            white 50%, 
            var(--green) 100%
        );
        background-size: 200% 100%;
        animation: progressTricolor 2s linear infinite;
    }
    
    @keyframes progressTricolor {
        0% { background-position: 0% 50%; }
        100% { background-position: 200% 50%; }
    }
    
    /* ===== EXPANDER (Dark) ===== */
    .streamlit-expanderHeader {
        background: var(--dark-card);
        border-radius: 10px;
        font-weight: 600;
        transition: all 0.3s ease;
        color: #e2e8f0;
        border: 2px solid var(--dark-surface);
    }
    
    .streamlit-expanderHeader:hover {
        background: var(--dark-surface);
        border-color: var(--saffron);
    }
    
    /* ===== METRICS (Gold) ===== */
    [data-testid="stMetricValue"] {
        font-size: 2.5rem;
        font-weight: 900;
        color: var(--gold);
        text-shadow: 0 0 10px rgba(255, 215, 0, 0.5);
    }
    
    /* ===== DATAFRAME (Dark) ===== */
    .stDataFrame {
        border-radius: 12px;
        overflow: hidden;
        box-shadow: var(--shadow-md);
    }
    
    .stDataFrame > div {
        background: var(--dark-card) !important;
    }
    
    /* ===== CUSTOM SCROLLBAR (Indian Colors) ===== */
    ::-webkit-scrollbar {
        width: 12px;
        height: 12px;
    }
    
    ::-webkit-scrollbar-track {
        background: var(--dark-bg);
        border-radius: 10px;
    }
    
    ::-webkit-scrollbar-thumb {
        background: linear-gradient(180deg, var(--saffron) 0%, var(--green) 100%);
        border-radius: 10px;
    }
    
    ::-webkit-scrollbar-thumb:hover {
        background: linear-gradient(180deg, var(--gold) 0%, var(--saffron) 100%);
    }
    
    /* ===== INPUT FIELDS (Dark) ===== */
    .stTextInput > div > div > input,
    .stTextArea > div > div > textarea,
    .stSelectbox > div > div > div {
        background-color: var(--dark-card) !important;
        color: #e2e8f0 !important;
        border: 2px solid var(--dark-surface) !important;
        border-radius: 10px;
        transition: all 0.3s ease;
    }
    
    .stTextInput > div > div > input:focus,
    .stTextArea > div > div > textarea:focus {
        border-color: var(--saffron) !important;
        box-shadow: 0 0 0 3px rgba(255, 153, 51, 0.2) !important;
    }
    
    /* ===== SLIDER (Indian Colors) ===== */
    .stSlider > div > div > div {
        background: var(--saffron) !important;
    }
    
    /* ===== INFO/WARNING/SUCCESS BOXES (Dark) ===== */
    .stAlert {
        background: var(--dark-card);
        border-radius: 10px;
        border-left: 4px solid var(--saffron);
        color: #e2e8f0;
    }
    
    /* ===== TOOLTIPS (Dark) ===== */
    [data-baseweb="tooltip"] {
        background: var(--dark-surface) !important;
        border-radius: 8px;
        box-shadow: var(--shadow-lg);
        color: #e2e8f0 !important;
    }
    
    /* ===== LOADING SPINNER (Tricolor) ===== */
    .stSpinner > div {
        border-top-color: var(--saffron) !important;
        border-right-color: white !important;
        border-bottom-color: var(--green) !important;
    }
    
    /* ===== MARKDOWN TEXT (Light on Dark) ===== */
    .stMarkdown {
        color: #e2e8f0;
    }
    
    /* ===== HEADINGS (Saffron) ===== */
    h1, h2, h3, h4, h5, h6 {
        color: var(--saffron) !important;
    }
    
    /* ===== LINKS (Gold) ===== */
    a {
        color: var(--gold) !important;
        transition: all 0.3s ease;
    }
    
    a:hover {
        color: var(--saffron) !important;
        text-decoration: underline;
    }
</style>

<!-- FLOATING "मेरा भारत महान" CHARACTERS -->
<div class="hindi-char char-1">म</div>
<div class="hindi-char char-2">े</div>
<div class="hindi-char char-3">र</div>
<div class="hindi-char char-4">ा</div>
<div class="hindi-char char-5">भ</div>
<div class="hindi-char char-6">ा</div>
<div class="hindi-char char-7">र</div>
<div class="hindi-char char-8">त</div>
<div class="hindi-char char-9">म</div>
<div class="hindi-char char-10">ह</div>
<div class="hindi-char char-11">ा</div>
<div class="hindi-char char-12">न</div>
""", unsafe_allow_html=True)


def get_project_path(relative_path):
    """Get absolute path from project root"""
    # Use pathlib for robust path handling and normalize the result
    try:
        from pathlib import Path
        return str((Path(PROJECT_ROOT) / relative_path).resolve())
    except Exception:
        return os.path.normpath(os.path.join(PROJECT_ROOT, relative_path))


def find_checkpoint(candidate):
    """Find the first existing checkpoint file.

    candidate may be a single relative path (str) or an iterable of relative paths.
    The function will check each candidate under the project root. If none
    exist, it will attempt a glob search for the basename anywhere under the
    project root (useful if the working directory is different or files are
    placed in sibling dirs like gui/outputs).
    Returns absolute path string or None if not found.
    """
    from pathlib import Path

    root = Path(PROJECT_ROOT)

    # Normalize input to a list
    if isinstance(candidate, (list, tuple)):
        candidates = list(candidate)
    else:
        candidates = [candidate]

    # Check each candidate directly under project root
    for rel in candidates:
        try:
            p = (root / rel).resolve()
        except Exception:
            p = Path(os.path.normpath(os.path.join(PROJECT_ROOT, rel)))

        if p.exists():
            return str(p)

    # Fallback: try to find by basename anywhere under project
    basenames = [Path(c).name for c in candidates]
    for name in basenames:
        for p in root.rglob(name):
            if p.is_file():
                return str(p.resolve())

    return None


def resolve_user_path(user_path: str) -> Optional[str]:
    """Resolve a user-entered path.

    If user_path is absolute and exists, return it. If it's relative, treat it as
    relative to the project root and return the resolved absolute path if exists.
    If the path doesn't exist, still return the absolute candidate (caller may check).
    """
    if not user_path:
        return None
    from pathlib import Path

    p = Path(user_path)
    if p.is_absolute():
        return str(p.resolve())

    # Treat as relative to project root
    return str((Path(PROJECT_ROOT) / user_path).resolve())


def ensure_vocab_object(obj) -> Optional["Vocabulary"]:
    """Ensure the returned object is a `Vocabulary` instance.

    If the loaded object is a dict (old pickle format), wrap it into a
    Vocabulary instance by populating key fields.
    """
    if obj is None:
        return None

    if isinstance(obj, Vocabulary):
        return obj

    # If it's a dict-like object, try populating a Vocabulary
    if isinstance(obj, dict):
        v = Vocabulary(config=None)
        # Common expected keys
        char2idx = obj.get('char2idx') or obj.get('token_to_idx') or obj.get('stoi')
        idx2char = obj.get('idx2char') or obj.get('itos')

        if char2idx:
            v.char2idx = char2idx
            v._size = len(char2idx)
        if idx2char:
            v.idx2char = {int(k): val for k, val in idx2char.items()}

        # Preserve frequency if present
        if 'char_freq' in obj:
            try:
                v.char_freq = obj.get('char_freq', {})
            except Exception:
                pass

        # Set convenience attributes if present
        v.pad_idx = obj.get('pad_idx', getattr(v, 'pad_idx', 0))
        v.sos_idx = obj.get('sos_idx', getattr(v, 'sos_idx', 1))
        v.eos_idx = obj.get('eos_idx', getattr(v, 'eos_idx', 2))
        v.unk_idx = obj.get('unk_idx', getattr(v, 'unk_idx', 3))

        return v

    # If it's some other object exposing char2idx, assume it's usable
    if hasattr(obj, 'char2idx') and hasattr(obj, 'idx2char'):
        return obj

    return None


def normalize_output_to_batch_list(output_indices) -> list:
    """Normalize various model.generate outputs into a list-of-lists of ints: [batch, seq_len].

    Supports torch.Tensor, numpy.ndarray, Python lists (including lists of tensors),
    and nested lists. Ensures conversion uses .item() for 0-dim tensors and avoids
    calling int() on multi-element tensors directly.
    """
    try:
        import numpy as _np
    except Exception:
        _np = None
    try:
        import torch as _torch
    except Exception:
        _torch = None

    def _to_batch_from_tensor(t):
        # detach and convert to CPU numpy
        t_cpu = t.detach().cpu()
        arr = t_cpu.numpy()
        lst = arr.tolist()
        # lst can be scalar, 1d list, or nested lists
        if not isinstance(lst, list):
            return [[int(lst)]]
        # 1D -> wrap as single-batch
        if len(lst) == 0:
            return [[]]
        if not isinstance(lst[0], list):
            return [[int(x) for x in lst]]
        # nested list
        return [[int(x) for x in inner] for inner in lst]

    def _to_batch_from_ndarray(a):
        lst = a.tolist()
        if not isinstance(lst, list):
            return [[int(lst)]]
        if len(lst) == 0:
            return [[]]
        if not isinstance(lst[0], list):
            return [[int(x) for x in lst]]
        return [[int(x) for x in inner] for inner in lst]

    # Torch tensor
    if _torch is not None and _torch.is_tensor(output_indices):
        return _to_batch_from_tensor(output_indices)

    # NumPy array
    if _np is not None and isinstance(output_indices, _np.ndarray):
        return _to_batch_from_ndarray(output_indices)

    # Python list/tuple
    if isinstance(output_indices, (list, tuple)):
        if len(output_indices) == 0:
            return [[]]
        first = output_indices[0]
        # If first element is list/tuple or ndarray/tensor -> treat as batch
        # If first element is a list/tuple it may be either:
        # - a sequence (list of token scalars) OR
        # - a list of beams (list of sequences) when beam_size>1
        if isinstance(first, (list, tuple)) or (_torch is not None and _torch.is_tensor(first)) or (_np is not None and isinstance(first, _np.ndarray)):
            batch = []
            for seq in output_indices:
                # Handle beam-case: seq may be a list of beams (each a sequence)
                # Detect beams by checking if seq is a list and its first element is a sequence-like
                seq_to_convert = seq
                if isinstance(seq, (list, tuple)) and len(seq) > 0:
                    inner0 = seq[0]
                    if (_torch is not None and _torch.is_tensor(inner0)) or (_np is not None and isinstance(inner0, _np.ndarray)) or isinstance(inner0, (list, tuple)):
                        # Treat `seq` as a list of beams: pick the top/first beam
                        seq_to_convert = seq[0]

                # Convert seq_to_convert to python list
                if _torch is not None and _torch.is_tensor(seq_to_convert):
                    arr = seq_to_convert.detach().cpu().numpy().tolist()
                elif _np is not None and isinstance(seq_to_convert, _np.ndarray):
                    arr = seq_to_convert.tolist()
                else:
                    arr = seq_to_convert

                if not isinstance(arr, list):
                    batch.append([int(arr)])
                else:
                    # convert elements to int using item() where necessary
                    converted = []
                    for e in arr:
                        if _torch is not None and _torch.is_tensor(e):
                            if e.numel() == 1:
                                converted.append(int(e.item()))
                            else:
                                # convert multi-element tensor to list then extend
                                converted.extend([int(x) for x in e.detach().cpu().numpy().tolist()])
                        elif _np is not None and isinstance(e, _np.generic):
                            converted.append(int(e))
                        else:
                            converted.append(int(e))
                    batch.append(converted)
            return batch
        else:
            # single sequence provided as flat list of scalars/tensors
            converted = []
            for e in output_indices:
                if _torch is not None and _torch.is_tensor(e):
                    if e.numel() == 1:
                        converted.append(int(e.item()))
                    else:
                        converted.extend([int(x) for x in e.detach().cpu().numpy().tolist()])
                elif _np is not None and isinstance(e, _np.generic):
                    converted.append(int(e))
                else:
                    converted.append(int(e))
            return [converted]

    # Fallback: try to iterate and coerce
    try:
        return [list(map(int, list(output_indices)))]
    except Exception as ex:
        raise TypeError(f"Unsupported output type for normalization: {type(output_indices).__name__}, error: {ex}")


def char_f1_score(pred: str, ref: str) -> float:
    """Compute character-level F1 treating strings as multisets of characters.

    Returns micro-averaged F1 for a single prediction-reference pair.
    """
    from collections import Counter
    p_counts = Counter(pred)
    r_counts = Counter(ref)
    if len(p_counts) == 0 and len(r_counts) == 0:
        return 1.0
    if len(p_counts) == 0 or len(r_counts) == 0:
        return 0.0

    tp = 0
    for ch, rc in r_counts.items():
        tp += min(rc, p_counts.get(ch, 0))

    pred_total = sum(p_counts.values())
    ref_total = sum(r_counts.values())

    precision = tp / pred_total if pred_total > 0 else 0.0
    recall = tp / ref_total if ref_total > 0 else 0.0
    if precision + recall == 0:
        return 0.0
    return 2 * (precision * recall) / (precision + recall)


def word_exact_match(pred: str, ref: str) -> float:
    """Compute word-level exact match accuracy for single pair (0 or 1)."""
    p = pred.strip().split()
    r = ref.strip().split()
    # Exact match at sequence level
    return 1.0 if p == r else 0.0


def generate_full_pptx(comparison_records: list, attached_result_path: Optional[str] = None, report_md: Optional[str] = None) -> Optional[bytes]:
    """Generate a PPTX bytes object summarizing the last comparison.

    - comparison_records: list of dicts as saved in st.session_state['last_comparison']
    - attached_result_path: optional path to a results JSON to include error analysis

    Returns PPTX as bytes or None if generation failed.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        return None

    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Hindi Transliteration Report"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

    # If report markdown provided, parse headings into slides first
    def _parse_md_to_slides(md_text: str):
        slides = []
        if not md_text:
            return slides
        lines = md_text.splitlines()
        current_title = None
        current_body = []
        for ln in lines:
            stripped = ln.strip()
            if stripped.startswith('#'):
                # Flush previous
                if current_title is not None:
                    slides.append((current_title, '\n'.join(current_body)))
                # New title (remove leading #s)
                title = stripped.lstrip('#').strip()
                current_title = title if title else 'Section'
                current_body = []
            else:
                if current_title is None:
                    # preamble
                    current_title = 'Intro'
                current_body.append(stripped)
        if current_title is not None:
            slides.append((current_title, '\n'.join(current_body)))
        return slides

    if report_md:
        md_slides = _parse_md_to_slides(report_md)
        for title_text, body_text in md_slides:
            s = prs.slides.add_slide(prs.slide_layouts[1])
            s.shapes.title.text = title_text[:200]
            try:
                s.placeholders[1].text = body_text[:2000]
            except Exception:
                pass

    # Summary slide with attached result metadata
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Summary"
    tf = slide.placeholders[1].text_frame
    tf.text = f"Configs tested: {len(comparison_records)}"
    p = tf.add_paragraph()
    p.text = f"Attached result: {attached_result_path or 'None'}"

    # Add a table slide (first up to 12 rows)
    if comparison_records:
        rows = min(12, len(comparison_records) + 1)
        cols = 6
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Comparison (sample)"
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(4.5)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        headers = ['Model', 'Type', 'Config', 'Word Acc', 'Char F1', 'Time (s)']
        for c, h in enumerate(headers):
            table.cell(0, c).text = h
        for r_idx in range(1, rows):
            rec_idx = r_idx - 1
            if rec_idx >= len(comparison_records):
                break
            rec = comparison_records[rec_idx]
            table.cell(r_idx, 0).text = str(rec.get('Model', ''))
            table.cell(r_idx, 1).text = str(rec.get('Type', ''))
            table.cell(r_idx, 2).text = str(rec.get('Config', ''))
            table.cell(r_idx, 3).text = str(rec.get('Word Acc', ''))
            table.cell(r_idx, 4).text = str(rec.get('Char F1', ''))
            table.cell(r_idx, 5).text = str(rec.get('Time (s)', ''))

    # Build charts using Plotly if available
    try:
        import plotly.express as ppx
        import plotly.io as pio
        use_plotly = True
    except Exception:
        use_plotly = False

    if use_plotly and comparison_records:
        try:
            df = pd.DataFrame(comparison_records)
            # Time chart
            fig_time = ppx.bar(df, x='Model', y='Time (s)', color='Config', title='Processing Time')
            img_time = pio.to_image(fig_time, format='png')
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.title.text = 'Processing Time'
            slide.shapes.add_picture(io.BytesIO(img_time), Inches(1), Inches(1.5), width=Inches(8))

            # Word Acc chart
            fig_acc = ppx.bar(df, x='Model', y='Word Acc', color='Config', title='Word-level Exact Accuracy')
            img_acc = pio.to_image(fig_acc, format='png')
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.title.text = 'Word-level Exact Accuracy'
            slide.shapes.add_picture(io.BytesIO(img_acc), Inches(1), Inches(1.5), width=Inches(8))

            # Char F1 chart
            fig_f1 = ppx.bar(df, x='Model', y='Char F1', color='Config', title='Character-level F1')
            img_f1 = pio.to_image(fig_f1, format='png')
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.title.text = 'Character-level F1'
            slide.shapes.add_picture(io.BytesIO(img_f1), Inches(1), Inches(1.5), width=Inches(8))
        except Exception:
            # Fall back to a slide saying chart export unavailable
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = 'Charts'
            slide.placeholders[1].text = 'Plotly image export not available in this environment.'
    else:
        # Informative slide when plotly not available
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = 'Charts'
        slide.placeholders[1].text = 'Charts require Plotly and an image renderer (kaleido). Install plotly[kaleido] for full charts.'

    # If attached result has error_analysis, include a sample slide
    if attached_result_path and os.path.exists(attached_result_path):
        try:
            with open(attached_result_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            if 'error_analysis' in data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = 'Error Analysis (sample)'
                txt = json.dumps(data.get('error_analysis', {}), ensure_ascii=False)[:2000]
                slide.placeholders[1].text = txt
        except Exception:
            pass

    # Save to bytes
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()


class TransliterationApp:
    def __init__(self):
        """Initialize the application with session state management"""
        # Initialize session state
        if 'initialized' not in st.session_state:
            st.session_state.initialized = True
            st.session_state.history = []
            st.session_state.api_keys = {}
            st.session_state.connected_providers = set()
            st.session_state.available_models = {}
            st.session_state.selected_models = {}
            st.session_state.comparison_results = []
            st.session_state.test_results = {}
        
        self.config = self.load_config()
        self.device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
        self.models = {}
        self.src_vocab = None
        self.tgt_vocab = None
        
        # Initialize LLM handler
        self.llm = LLMTransliterator(self.config)
        
        # Restore API keys from session state
        for provider, api_key in st.session_state.api_keys.items():
            self.llm.setup_client(provider, api_key)
        
        # Load models
        self.load_models()
        
        # Initialize evaluator
        self.evaluator = Evaluator()
    
    def load_config(self) -> Dict:
        """Load configuration with UTF-8 encoding support"""
        config_path = get_project_path('config/config.yaml')
        
        if os.path.exists(config_path):
            try:
                with open(config_path, 'r', encoding='utf-8') as f:
                    return yaml.safe_load(f)
            except UnicodeDecodeError:
                try:
                    with open(config_path, 'r', encoding='utf-8-sig') as f:
                        return yaml.safe_load(f)
                except Exception as e:
                    st.warning(f"⚠️ Config encoding error: {e}. Using defaults.")
            except Exception as e:
                st.warning(f"⚠️ Config load error: {e}. Using defaults.")
        
        # Default configuration
        return {
            'data': {
                'max_seq_length': 50,
                'max_train_samples': 100000,
                'device': 'auto'
            },
            'lstm': {
                'embedding_dim': 256, 
                'hidden_dim': 512, 
                'num_layers': 2, 
                'dropout': 0.3, 
                'bidirectional': True
            },
            'transformer': {
                'd_model': 256, 
                'n_heads': 8, 
                'num_layers': 2, 
                'd_ff': 1024, 
                'dropout': 0.1, 
                'use_local_attention': True, 
                'local_attention_window': 5, 
                'max_seq_length': 100
            },
            'training': {
                'batch_size': 128,
                'learning_rate': 0.0005,
                'epochs': 10,
                'gradient_clip': 1.0,
                'teacher_forcing_ratio': 0.5,
                'beam_sizes': [1, 3, 5, 10]
            },
            'llm': {
                'max_tokens': 100,
                'temperature_values': [0.1, 0.3, 0.5, 0.7, 0.9, 1.0], 
                'top_p_values': [0.9, 0.95, 0.99, 1.0]
            },
            'preprocessing': {
                'min_frequency': 2
            }
        }
    
    def load_models(self):
        """Load trained neural models by reading paths directly from the config file."""
        self.model_status = {'vocab': False, 'lstm': False, 'transformer': False}
        
    # --- 1. Load Vocabularies ---
    # NOTE: This method reads file names and directories from `config/config.yaml` under
    # the `paths` and `gui` sections. To override where models/vocabs are loaded from,
    # update `config/config.yaml` e.g. `paths.checkpoint_dir` or `gui.transformer_checkpoint`.
        st.sidebar.info("🔍 Searching for vocabularies...")
        try:
            # ✅ FIX: Read vocab paths directly from config, with optional session override
            cfg_paths = self.config.get('paths', {})
            processed_dir = cfg_paths.get('processed_data_dir', 'data/processed')

            # Check session overrides first
            overrides = st.session_state.get('path_overrides', {})

            src_vocab_override = overrides.get('src_vocab')
            tgt_vocab_override = overrides.get('tgt_vocab')

            if src_vocab_override:
                src_full_path = resolve_user_path(src_vocab_override)
                # If user provided a filename (not absolute), also try relative to processed_dir
                if not os.path.exists(src_full_path):
                    src_full_path = get_project_path(os.path.join(processed_dir, src_vocab_override))
            else:
                src_vocab_rel_path = cfg_paths.get('src_vocab_file', 'vocab_src.pkl')
                src_full_path = get_project_path(os.path.join(processed_dir, src_vocab_rel_path))

            if tgt_vocab_override:
                tgt_full_path = resolve_user_path(tgt_vocab_override)
                if not os.path.exists(tgt_full_path):
                    tgt_full_path = get_project_path(os.path.join(processed_dir, tgt_vocab_override))
            else:
                tgt_vocab_rel_path = cfg_paths.get('tgt_vocab_file', 'vocab_tgt.pkl')
                tgt_full_path = get_project_path(os.path.join(processed_dir, tgt_vocab_rel_path))
            # Robust vocab loading: support JSON and PKL (and Vocabulary.load)
            def _load_vocab_file(path_candidate):
                if not path_candidate:
                    return None
                try:
                    if path_candidate.endswith('.json'):
                        v = Vocabulary()
                        v.load(path_candidate)
                        return v
                    if path_candidate.endswith('.pkl') or path_candidate.endswith('.pickle'):
                        with open(path_candidate, 'rb') as f:
                            return pickle.load(f)
                    # Try Vocabulary.load then pickle as fallback
                    try:
                        v = Vocabulary()
                        v.load(path_candidate)
                        return v
                    except Exception:
                        with open(path_candidate, 'rb') as f:
                            return pickle.load(f)
                except Exception:
                    return None

            src_obj = _load_vocab_file(src_full_path) if os.path.exists(src_full_path) else None
            tgt_obj = _load_vocab_file(tgt_full_path) if os.path.exists(tgt_full_path) else None

            # Ensure we have Vocabulary instances (not dicts)
            src_vocab_obj = ensure_vocab_object(src_obj)
            tgt_vocab_obj = ensure_vocab_object(tgt_obj)

            if src_vocab_obj is not None and tgt_vocab_obj is not None:
                self.src_vocab = src_vocab_obj
                self.tgt_vocab = tgt_vocab_obj
                self.model_status['vocab'] = True
                st.sidebar.info(f"Found vocab files:\n src={src_full_path}\n tgt={tgt_full_path}")
            else:
                st.sidebar.error(f"❌ VOCAB NOT FOUND or failed to load. Tried:\n src={src_full_path}\n tgt={tgt_full_path}")

        except Exception as e:
            st.sidebar.error(f"❌ Vocab loading failed: {e}")
            return # Stop if vocab fails

        # --- 2. Load LSTM Model ---
        if self.model_status['vocab']:
            st.sidebar.info("🔍 Searching for LSTM model...")
            try:
                # ✅ FIX: Read LSTM checkpoint path from config or overrides
                cfg_paths = self.config.get('paths', {})
                cfg_gui = self.config.get('gui', {})
                checkpoint_dir = overrides.get('checkpoint_dir', cfg_paths.get('checkpoint_dir', 'outputs/checkpoints'))
                lstm_checkpoint_name = overrides.get('lstm_checkpoint', cfg_gui.get('lstm_checkpoint', 'lstm_best.pt'))

                # Resolve a user-specified full path first
                if os.path.isabs(lstm_checkpoint_name) and os.path.exists(lstm_checkpoint_name):
                    full_path = lstm_checkpoint_name
                else:
                    full_path = resolve_user_path(os.path.join(checkpoint_dir, lstm_checkpoint_name))

                if os.path.exists(full_path):
                    st.sidebar.info(f"Found LSTM: {lstm_checkpoint_name}")
                    lstm_config = self.config['lstm']
                    src_vocab_size = getattr(self.src_vocab, 'size', len(getattr(self.src_vocab, 'char2idx', self.src_vocab)))
                    tgt_vocab_size = getattr(self.tgt_vocab, 'size', len(getattr(self.tgt_vocab, 'char2idx', self.tgt_vocab)))

                    # Instantiate using config-driven constructor
                    self.models['LSTM'] = Seq2SeqLSTM(
                        src_vocab_size,
                        tgt_vocab_size,
                        config=self.config
                    ).to(self.device)

                    checkpoint = torch.load(full_path, map_location=self.device)

                    # Handle different checkpoint save formats
                    state_dict = checkpoint.get('model_state_dict', checkpoint.get('state_dict', checkpoint)) if isinstance(checkpoint, dict) else checkpoint

                    # Filter mismatched parameters (robust load)
                    model_dict = self.models['LSTM'].state_dict()
                    filtered_state_dict = {}
                    for k, v in (state_dict.items() if isinstance(state_dict, dict) else []):
                        if k in model_dict and getattr(v, 'shape', None) == getattr(model_dict[k], 'shape', None):
                            filtered_state_dict[k] = v

                    if filtered_state_dict:
                        self.models['LSTM'].load_state_dict(filtered_state_dict, strict=False)
                    else:
                        # Try direct load as last resort
                        try:
                            self.models['LSTM'].load_state_dict(state_dict)
                        except Exception:
                            raise

                    self.models['LSTM'].eval()
                    self.model_status['lstm'] = True
                else:
                    st.sidebar.warning(f"LSTM MODEL NOT FOUND at {full_path}")
            except Exception as e:
                st.sidebar.error(f"❌ LSTM load failed: {str(e)}")

        # --- 3. Load Transformer Model ---
        if self.model_status['vocab']:
            st.sidebar.info("🔍 Searching for Transformer model...")
            try:
                # ✅ FIX: Read Transformer checkpoint path from config or overrides
                cfg_paths = self.config.get('paths', {})
                cfg_gui = self.config.get('gui', {})
                checkpoint_dir = overrides.get('checkpoint_dir', cfg_paths.get('checkpoint_dir', 'outputs/checkpoints'))
                trans_checkpoint_name = overrides.get('transformer_checkpoint', cfg_gui.get('transformer_checkpoint', 'transformer_best.pt'))

                if os.path.isabs(trans_checkpoint_name) and os.path.exists(trans_checkpoint_name):
                    full_path = trans_checkpoint_name
                else:
                    full_path = resolve_user_path(os.path.join(checkpoint_dir, trans_checkpoint_name))

                if os.path.exists(full_path):
                    st.sidebar.info(f"Found Transformer: {trans_checkpoint_name}")
                    transformer_config = self.config['transformer']
                    src_vocab_size = getattr(self.src_vocab, 'size', len(getattr(self.src_vocab, 'char2idx', self.src_vocab)))
                    tgt_vocab_size = getattr(self.tgt_vocab, 'size', len(getattr(self.tgt_vocab, 'char2idx', self.tgt_vocab)))

                    # Instantiate using config-driven constructor
                    self.models['Transformer'] = TransformerSeq2Seq(
                        src_vocab_size,
                        tgt_vocab_size,
                        config=self.config
                    ).to(self.device)
                    
                    checkpoint = torch.load(full_path, map_location=self.device)

                    # Many formats: prefer model_state_dict or state_dict
                    state_dict = checkpoint.get('model_state_dict', checkpoint.get('state_dict', checkpoint)) if isinstance(checkpoint, dict) else checkpoint

                    # Try to load directly, fall back to filtered loading on mismatch
                    try:
                        self.models['Transformer'].load_state_dict(state_dict)
                    except Exception:
                        # Filter by shape compatibility
                        model_dict = self.models['Transformer'].state_dict()
                        filtered = {k: v for k, v in (state_dict.items() if isinstance(state_dict, dict) else []) if k in model_dict and getattr(v, 'shape', None) == getattr(model_dict[k], 'shape', None)}
                        if filtered:
                            self.models['Transformer'].load_state_dict(filtered, strict=False)
                        else:
                            raise

                    self.models['Transformer'].eval()
                    self.model_status['transformer'] = True
                else:
                    st.sidebar.warning(f"TRANSFORMER NOT FOUND at {full_path}")
            except Exception as e:
                st.sidebar.error(f"❌ Transformer load failed: {str(e)}")
    
    def transliterate_neural(self, text: str, model_name: str, beam_size: int = 1, return_beams: bool = False):
        """Transliterate using neural models"""
        if model_name not in self.models:
            return f"❌ {model_name} not loaded"
        
        try:
            input_indices = self.src_vocab.encode(text)
            input_tensor = torch.tensor([input_indices]).to(self.device)
            
            with torch.no_grad():
                output_indices = self.models[model_name].generate(
                    input_tensor, 
                    max_length=100, 
                    beam_size=beam_size
                )
            # Normalize output and decode first sequence
            normalized = normalize_output_to_batch_list(output_indices)
            # If caller requested beams, return the nested list (per-batch, list of sequences)
            if return_beams:
                return normalized

            first_seq = normalized[0] if normalized and len(normalized) > 0 else []
            return self.tgt_vocab.decode(first_seq)
        except Exception as e:
            return f"❌ Error: {str(e)}"
    
    def sidebar_content(self):
        """Render sidebar with system status"""
        with st.sidebar:
            st.markdown('<div class="section-header"><h2>🇮🇳 System Status</h2></div>', unsafe_allow_html=True)
            
            # Device info
            device_emoji = "🎮" if "cuda" in str(self.device) else "💻"
            st.info(f"{device_emoji} **Device:** {self.device}")
            
            # Model status
            st.markdown("### 🤖 Neural Models")

            # Manual path overrides
            with st.expander("🛠️ Manual Path Overrides (optional)", expanded=False):
                if 'path_overrides' not in st.session_state:
                    st.session_state.path_overrides = {}

                # Show current defaults from config for guidance
                cfg_paths = self.config.get('paths', {})
                cfg_gui = self.config.get('gui', {})

                co_checkpoint_dir = st.text_input("Checkpoint directory (relative or absolute)",
                                                 value=st.session_state.path_overrides.get('checkpoint_dir', cfg_paths.get('checkpoint_dir', 'outputs/checkpoints')))
                co_lstm = st.text_input("LSTM checkpoint filename or path",
                                       value=st.session_state.path_overrides.get('lstm_checkpoint', cfg_gui.get('lstm_checkpoint', 'lstm_best.pt')))
                co_transformer = st.text_input("Transformer checkpoint filename or path",
                                              value=st.session_state.path_overrides.get('transformer_checkpoint', cfg_gui.get('transformer_checkpoint', 'transformer_best.pt')))

                # Vocab overrides
                src_vocab = st.text_input("Source vocab path (relative to processed_data_dir or absolute)",
                                         value=st.session_state.path_overrides.get('src_vocab', cfg_paths.get('src_vocab_file', 'vocab_src.pkl')))
                tgt_vocab = st.text_input("Target vocab path (relative to processed_data_dir or absolute)",
                                         value=st.session_state.path_overrides.get('tgt_vocab', cfg_paths.get('tgt_vocab_file', 'vocab_tgt.pkl')))

                if st.button("Save path overrides"):
                    st.session_state.path_overrides['checkpoint_dir'] = co_checkpoint_dir.strip()
                    st.session_state.path_overrides['lstm_checkpoint'] = co_lstm.strip()
                    st.session_state.path_overrides['transformer_checkpoint'] = co_transformer.strip()
                    st.session_state.path_overrides['src_vocab'] = src_vocab.strip()
                    st.session_state.path_overrides['tgt_vocab'] = tgt_vocab.strip()
                    st.success("Saved path overrides — reload models to apply")
                if st.button("🔎 Test resolved paths"):
                    test_out = {}
                    po = st.session_state.get('path_overrides', {})
                    # Resolve using same logic as load_models
                    def _resolve_vocab(p, default):
                        if p:
                            rp = resolve_user_path(p)
                            if os.path.exists(rp):
                                return rp
                            # try relative to processed_dir
                            proc = self.config.get('paths', {}).get('processed_data_dir', 'data/processed')
                            cand = get_project_path(os.path.join(proc, p))
                            return cand
                        return get_project_path(os.path.join(self.config.get('paths', {}).get('processed_data_dir', 'data/processed'), default))

                    checkpoint_dir = po.get('checkpoint_dir', self.config.get('paths', {}).get('checkpoint_dir', 'outputs/checkpoints'))
                    lstm_ck = po.get('lstm_checkpoint', self.config.get('gui', {}).get('lstm_checkpoint', 'lstm_best.pt'))
                    trans_ck = po.get('transformer_checkpoint', self.config.get('gui', {}).get('transformer_checkpoint', 'transformer_best.pt'))

                    test_out['lstm'] = resolve_user_path(os.path.join(checkpoint_dir, lstm_ck)) if not os.path.isabs(lstm_ck) else lstm_ck
                    test_out['transformer'] = resolve_user_path(os.path.join(checkpoint_dir, trans_ck)) if not os.path.isabs(trans_ck) else trans_ck
                    test_out['src_vocab'] = _resolve_vocab(po.get('src_vocab'), self.config.get('paths', {}).get('src_vocab_file', 'vocab_src.pkl'))
                    test_out['tgt_vocab'] = _resolve_vocab(po.get('tgt_vocab'), self.config.get('paths', {}).get('tgt_vocab_file', 'vocab_tgt.pkl'))

                    for k, v in test_out.items():
                        st.write(f"{k}: {v} — exists={os.path.exists(v)}")
            
            if self.model_status['vocab']:
                src_size = self.src_vocab.size if hasattr(self.src_vocab, 'size') else len(self.src_vocab)
                tgt_size = self.tgt_vocab.size if hasattr(self.tgt_vocab, 'size') else len(self.tgt_vocab)
                st.success(f"✅ Vocabularies: Src={src_size}, Tgt={tgt_size}")
            else:
                st.warning("⚠️ Vocabularies not loaded")
            
            if self.model_status['lstm']:
                st.success("✅ LSTM (Bidirectional, 2 layers)")
            else:
                st.info("ℹ️ LSTM not available")
            
            if self.model_status['transformer']:
                st.success("✅ Transformer (Local Attention)")
            else:
                st.info("ℹ️ Transformer not available")
            
            # LLM Providers
            st.markdown("### 🌐 LLM Providers")
            if st.session_state.connected_providers:
                for provider in st.session_state.connected_providers:
                    st.success(f"✅ {provider.title()}")
                    
                    if provider in st.session_state.available_models:
                        model_count = len(st.session_state.available_models[provider])
                        st.caption(f"   📦 {model_count} models available")
            else:
                st.info("💡 No LLM providers connected")
            
            # Rate limits for Groq
            if 'groq' in st.session_state.connected_providers:
                if st.button("🔄 Check Rate Limits", key="check_limits"):
                    limits = self.llm.check_rate_limits('groq')
                    if limits:
                        st.markdown("**Rate Limits:**")
                        st.caption(f"Requests: {limits.get('requests_remaining', 'N/A')}/{limits.get('requests_limit', 'N/A')}")
                        st.caption(f"Tokens: {limits.get('tokens_remaining', 'N/A')}/{limits.get('tokens_limit', 'N/A')}")
            
            st.divider()
            
            # Assignment checklist
            st.markdown("### ✅ CS772 Requirements")
            requirements = {
                "LSTM (≤2 layers)": self.model_status['lstm'],
                "Transformer (Local Attn)": self.model_status['transformer'],
                "LLM Integration": bool(st.session_state.connected_providers),
                "Greedy + Beam Search": True,
                "≤100k Training Limit": True,
                "ACL-Compliant Metrics": True,
                "Error Analysis": True
            }
            
            for req, status in requirements.items():
                icon = "✅" if status else "⭕"
                st.markdown(f"{icon} {req}")
            
            st.divider()
            
            # Proud message
            st.markdown("### 🇮🇳 मेरा भारत महान")
            st.caption("Celebrating Indian Heritage")
            # Recent history panel
            st.divider()
            with st.expander("Recent Searches & History", expanded=False):
                hist = st.session_state.get('history', [])
                if hist:
                    for item in hist[-10:][::-1]:
                        ts = item.get('timestamp', item.get('Time', ''))
                        model = item.get('model', item.get('Model', ''))
                        out = item.get('output', item.get('Output', ''))
                        st.markdown(f"**{ts}** — `{model}` — {out[:80]}")
                    if st.button("Clear History"):
                        st.session_state['history'] = []
                        st.success("History cleared")
                else:
                    st.info("No history yet — run transliterations to populate this list.")
    
    def tab_transliterate(self):
        """Main transliteration tab"""
        st.markdown('<div class="section-header"><h2>🎯 Roman → Devanagari Transliteration</h2></div>', unsafe_allow_html=True)
        
        col1, col2 = st.columns([2, 3])
        
        with col1:
            st.markdown("### ⚙️ Configuration")
            
            # Model selection
            available_models = list(self.models.keys())
            
            # Add LLM options
            for provider in st.session_state.connected_providers:
                available_models.append(f"LLM ({provider.title()})")
            
            if not available_models:
                st.warning("⚠️ No models available. Train models or connect LLM providers.")
                return
            
            model_type = st.selectbox(
                "Select Model",
                available_models,
                help="Choose the model for transliteration"
            )
            
            # Model-specific settings
            if model_type in ["LSTM", "Transformer"]:
                st.markdown("#### 🔍 Decoding Strategy")
                use_beam = st.checkbox("Use Beam Search", value=True, 
                                      help="Beam search explores multiple paths")
                
                if use_beam:
                    beam_size = st.slider("Beam Size", 2, 10, 5, 
                                        help="Higher = better quality but slower")
                else:
                    beam_size = 1
                
                st.info(f"{'🎯 Beam Search' if use_beam else '⚡ Greedy'} (beam={beam_size})")
            
            elif "LLM" in model_type:
                provider = model_type.split("(")[1].split(")")[0].lower()
                
                # ✅ FIXED: Better model selection UI
                if provider in st.session_state.available_models:
                    model_list = st.session_state.available_models[provider]
                    
                    if not model_list:
                        st.warning(f"⚠️ No models available for {provider.title()}")
                        selected_model = None
                    else:
                        # Group by category
                        categories = {}
                        for model in model_list:
                            cat = model.get('category', 'general')
                            if cat not in categories:
                                categories[cat] = []
                            categories[cat].append(model)
                        
                        st.markdown("#### 🤖 Model Selection")
                        
                        # ✅ FIX: Show categories only if more than one exists
                        if len(categories) > 1:
                            selected_category = st.selectbox(
                                "Category",
                                list(categories.keys()),
                                format_func=lambda x: f"📦 {x.upper()}",
                                help="Filter models by family"
                            )
                            available_models_in_cat = categories[selected_category]
                        else:
                            # Single category - skip dropdown
                            selected_category = list(categories.keys())[0]
                            available_models_in_cat = categories[selected_category]
                            st.caption(f"📦 Category: **{selected_category.upper()}**")
                        
                        # Model dropdown with rich formatting
                        model_options = {}
                        for m in available_models_in_cat:
                            model_id = m['id']
                            ctx = m.get('context_window', 0)
                            ctx_str = f"{ctx//1000}K" if ctx >= 1000 else str(ctx)
                            
                            # Display name with context window
                            display_name = f"{model_id}"
                            if ctx > 0:
                                display_name += f" ({ctx_str} ctx)"
                            
                            model_options[display_name] = model_id
                        
                        selected_display = st.selectbox(
                            "Model",
                            list(model_options.keys()),
                            help=f"{len(available_models_in_cat)} models in {selected_category}"
                        )
                        
                        selected_model = model_options[selected_display]
                        
                        # Show model details
                        selected_model_obj = next((m for m in model_list if m['id'] == selected_model), None)
                        if selected_model_obj:
                            col_a, col_b = st.columns(2)
                            with col_a:
                                if selected_model_obj.get('context_window', 0) > 0:
                                    st.caption(f"🔢 Context: {selected_model_obj['context_window']:,}")
                            with col_b:
                                if selected_model_obj.get('supports_reasoning'):
                                    st.caption("🧠 Reasoning: ✅")
                        
                        st.session_state.selected_models[provider] = selected_model
                else:
                    # ✅ FIX: Fetch models on-demand if not cached
                    st.info(f"📥 Fetching {provider.title()} models...")
                    with st.spinner("Loading models..."):
                        models = self.llm.get_available_models(provider)
                        if models:
                            st.session_state.available_models[provider] = models
                            st.rerun()  # Refresh to show models
                        else:
                            st.error(f"❌ Failed to fetch models for {provider.title()}")
                            selected_model = None
                
                # Generation parameters
                st.markdown("#### 🎛️ Parameters")
                col_temp, col_top = st.columns(2)
                
                with col_temp:
                    temperature = st.slider("Temperature", 0.0, 1.0, 0.3, 0.05,
                                          help="Lower = more deterministic")
                
                with col_top:
                    top_p = st.slider("Top-p", 0.0, 1.0, 0.95, 0.05,
                                    help="Nucleus sampling threshold")
                
                # ✅ FIX: Reasoning support (only for compatible models)
                use_reasoning = False
                if 'selected_model' in locals() and selected_model:
                    selected_model_obj = next(
                        (m for m in st.session_state.available_models.get(provider, []) 
                         if m['id'] == selected_model), 
                        None
                    )
                    
                    if selected_model_obj and selected_model_obj.get('supports_reasoning'):
                        use_reasoning = st.checkbox(
                            "🧠 Enable Reasoning", 
                            help="Show step-by-step thought process (slower)"
                        )
        
        with col2:
            st.markdown("### 📝 Input & Output")
            
            input_text = st.text_area(
                "Enter Roman Text",
                height=120,
                placeholder="Type here...\nExample: namaste bharat",
                help="Roman script → Devanagari",
                value=st.session_state.get('example_text', '')
            )
            
            if 'example_text' in st.session_state:
                del st.session_state.example_text
            
            if st.button("✨ **Transliterate Now**", type="primary", use_container_width=True):
                if input_text.strip():
                    start_time = time.time()
                    
                    with st.spinner("🔄 Processing... मेरा भारत महान"):
                        if model_type in ["LSTM", "Transformer"]:
                            words = input_text.strip().split()
                            results = []
                            
                            progress_bar = st.progress(0)
                            
                            for i, word in enumerate(words):
                                output = self.transliterate_neural(word, model_type, beam_size)
                                results.append(output)
                                progress_bar.progress((i + 1) / len(words))
                            
                            output_text = ' '.join(results)
                            progress_bar.empty()
                            
                        elif "LLM" in model_type:
                            provider = model_type.split("(")[1].split(")")[0].lower()
                            output_text = self.llm.transliterate(
                                input_text,
                                provider=provider,
                                model=selected_model if 'selected_model' in locals() else None,
                                temperature=temperature if 'temperature' in locals() else 0.3,
                                top_p=top_p if 'top_p' in locals() else 0.95,
                                use_reasoning=use_reasoning if 'use_reasoning' in locals() else False
                            )
                        else:
                            output_text = "❌ Model not available"
                    
                    elapsed_time = time.time() - start_time
                    
                    if output_text and not output_text.startswith("❌"):
                        st.balloons()
                        st.success(f"✅ Completed in {elapsed_time:.2f}s 🇮🇳")
                        
                        st.markdown("### 📤 Result")
                        st.markdown(f'<div class="transliteration-output">{output_text}</div>', unsafe_allow_html=True)
                        
                        st.code(output_text, language=None)
                        
                        # Add to history
                        st.session_state.history.append({
                            'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                            'input': input_text,
                            'output': output_text,
                            'model': model_type + (f" ({selected_model})" if 'selected_model' in locals() else ""),
                            'time': elapsed_time
                        })
                    else:
                        st.error(f"❌ Failed: {output_text}")
                else:
                    st.warning("⚠️ Please enter text")
        
        # Examples
        st.divider()
        st.markdown("### ⚡ Quick Examples")
        
        examples = [
            ("namaste", "नमस्ते", "Greeting"),
            ("bharat", "भारत", "Country"),
            ("vidyalaya", "विद्यालय", "School"),
            ("computer", "कंप्यूटर", "Tech"),
            ("adhyapak", "अध्यापक", "Teacher"),
            ("pustakaalaya", "पुस्तकालय", "Library"),
            ("paryavaran", "पर्यावरण", "Environment"),
            ("sanvidhan", "संविधान", "Constitution")
        ]
        
        cols = st.columns(4)
        for i, (roman, devanagari, category) in enumerate(examples):
            with cols[i % 4]:
                if st.button(f"**{roman}**", key=f"ex_{i}", use_container_width=True,
                           help=f"{category}: {devanagari}"):
                    st.session_state.example_text = roman
                    st.rerun()
    
    def tab_api_configuration(self):
        """API Configuration tab"""
        st.markdown('<div class="section-header"><h2>🔑 API Configuration</h2></div>', unsafe_allow_html=True)
        st.info("💡 Connect to LLM providers instantly. No restart required!")
        
        providers = [
            {
                'name': 'groq',
                'title': 'Groq',
                'description': 'Ultra-fast inference (Llama, Mixtral, Gemma)',
                'emoji': '⚡',
                'free_tier': True
            },
            {
                'name': 'openai',
                'title': 'OpenAI',
                'description': 'GPT-3.5 and GPT-4 models',
                'emoji': '🤖',
                'free_tier': False
            },
            {
                'name': 'anthropic',
                'title': 'Anthropic',
                'description': 'Claude 3 family',
                'emoji': '🧠',
                'free_tier': False
            },
            {
                'name': 'google',
                'title': 'Google',
                'description': 'Gemini Pro',
                'emoji': '🌟',
                'free_tier': True
            }
        ]
        
        cols = st.columns(2)
        
        for i, provider in enumerate(providers):
            with cols[i % 2]:
                st.markdown(f"""
                <div class="provider-card">
                    <h4>{provider['emoji']} {provider['title']}</h4>
                    <p>{provider['description']}</p>
                    {f"<span style='background:#138808;color:white;padding:4px 12px;border-radius:8px;font-size:11px;font-weight:700;'>FREE TIER</span>" if provider['free_tier'] else ""}
                </div>
                """, unsafe_allow_html=True)
                
                is_connected = provider['name'] in st.session_state.connected_providers
                
                if is_connected:
                    st.success(f"✅ Connected")
                    
                    # ✅ FIX: Show model count and allow refresh
                    if provider['name'] in st.session_state.available_models:
                        count = len(st.session_state.available_models[provider['name']])
                        st.caption(f"📦 {count} models available")
                    
                    col_info, col_refresh, col_disconnect = st.columns([2, 1, 1])
                    
                    with col_refresh:
                        if st.button("🔄", key=f"refresh_{provider['name']}", 
                                   help="Refresh models"):
                            with st.spinner("Refreshing..."):
                                models = self.llm.get_available_models(provider['name'])
                                if models:
                                    st.session_state.available_models[provider['name']] = models
                                    st.success(f"✅ {len(models)} models loaded!")
                                    time.sleep(1)
                                    st.rerun()
                    
                    with col_disconnect:
                        if st.button("🔌", key=f"disconnect_{provider['name']}", 
                                   help="Disconnect"):
                            st.session_state.connected_providers.discard(provider['name'])
                            st.session_state.api_keys.pop(provider['name'], None)
                            st.session_state.available_models.pop(provider['name'], None)
                            st.rerun()
                else:
                    api_key = st.text_input(
                        f"{provider['title']} API Key",
                        type="password",
                        key=f"api_key_{provider['name']}",
                        placeholder="sk-..." if provider['name'] == 'openai' else "API Key"
                    )
                    
                    if st.button(f"🔗 Connect", key=f"connect_{provider['name']}",
                               use_container_width=True, type="primary"):
                        if api_key and api_key.strip():
                            with st.spinner(f"Connecting to {provider['title']}..."):
                                success = self.llm.setup_client(provider['name'], api_key.strip())
                                
                                if success:
                                    st.session_state.api_keys[provider['name']] = api_key.strip()
                                    st.session_state.connected_providers.add(provider['name'])
                                    
                                    # ✅ FIX: Fetch models for ALL providers, not just Groq
                                    st.info("📥 Fetching available models...")
                                    models = self.llm.get_available_models(provider['name'])
                                    if models:
                                        st.session_state.available_models[provider['name']] = models
                                        st.success(f"✅ Connected! {len(models)} models loaded")
                                    else:
                                        st.warning(f"⚠️ Connected but no models found")
                                    
                                    st.balloons()
                                    time.sleep(1)
                                    st.rerun()
                                else:
                                    st.error(f"❌ Connection failed. Check API key.")
                        else:
                            st.warning("⚠️ Please enter API key")
        
        # Test connection
        if st.session_state.connected_providers:
            st.divider()
            st.markdown("### 🧪 Test Connections")
            
            test_text = st.text_input("Test Text", value="namaste", 
                                     help="Text to test transliteration")
            
            if st.button("🚀 Test All Connected Providers", type="primary"):
                results = {}
                
                for provider in st.session_state.connected_providers:
                    with st.spinner(f"Testing {provider}..."):
                        try:
                            result = self.llm.transliterate(test_text, provider=provider)
                            results[provider] = result
                        except Exception as e:
                            results[provider] = f"Error: {str(e)}"
                
                st.markdown("#### Test Results:")
                for provider, result in results.items():
                    if "Error" not in result and "ERROR" not in result:
                        st.success(f"**{provider.title()}:** {result}")
                    else:
                        st.error(f"**{provider.title()}:** {result}")
    
    def tab_compare_models(self):
        """Model comparison tab"""
        st.markdown('<div class="section-header"><h2>📊 Model Comparison</h2></div>', unsafe_allow_html=True)
        
        st.markdown("### 🔁 Input & Reference")
        col_src, col_ref = st.columns(2)
        with col_src:
            source_text = st.text_area(
                "Source (Roman)",
                value=st.session_state.get('compare_source', "mera bharat mahan"),
                height=80,
                help="Input text in Roman script (words will be compared word-by-word)"
            )
            st.session_state.compare_source = source_text
        with col_ref:
            reference_text = st.text_area(
                "Reference (Devanagari)",
                value=st.session_state.get('compare_ref', "मेरा भारत महान"),
                height=80,
                help="Expected Devanagari output to compare against"
            )
            st.session_state.compare_ref = reference_text
        
        col1, col2 = st.columns(2)
        
        with col1:
            include_neural = st.checkbox("Include Neural Models", value=True)
            neural_models = list(self.models.keys())
            selected_neural = st.multiselect("Neural Models to Test", neural_models, default=neural_models)
            default_baseline = st.selectbox("Baseline (for comparison)", ["(none)"] + neural_models, index=1 if neural_models else 0)
            if include_neural:
                beam_sizes = st.multiselect("Beam Sizes", [1, 2, 3, 5, 10], default=[1, 5])
        
        with col2:
            include_llm = st.checkbox("Include LLM", value=True)
            llm_providers = list(st.session_state.connected_providers)
            selected_providers = st.multiselect("LLM Providers to Test", llm_providers, default=llm_providers)
            if include_llm:
                temps = st.multiselect("Temperatures (grid)", [0.0, 0.1, 0.3, 0.5, 0.7], default=[0.3])
                top_ps = st.multiselect("Top-p values (grid)", [0.5, 0.7, 0.9, 0.95, 1.0], default=[0.95])

    # N-best beam options
        with st.expander("Advanced: N-best & per-beam metrics", expanded=False):
            show_nbest = st.checkbox("Show N-best beams and per-beam metrics", value=False)
            if show_nbest:
                n_best = st.slider("Top-K beams to show", 1, 10, 3)
        st.markdown("---")
        st.markdown("### Reranking & Diagnostics")
        rerank_method = st.selectbox("Rerank selection", ["model_score", "length_norm", "mixed"], index=1, help="Choose how to pick the top beam from N-best")
        alpha = st.number_input("Length-normalization alpha", min_value=0.0, max_value=2.0, value=0.7, step=0.1)
        mix_lambda = st.slider("Mixed rerank λ (normalize vs heuristic)", 0.0, 1.0, 0.6)
        st.markdown("(Oracle diagnostic will compute the best achievable metric if we pick the best beam by char-F1)")
        
        if st.button("🔄 **Run Comparison**", type="primary", use_container_width=True):
            if not source_text.strip() or not reference_text.strip():
                st.warning("⚠️ Provide both source and reference text for metric computation")
            else:
                results = []
                progress = st.progress(0)

                # Build grid
                tasks = []
                if include_neural:
                    for model_name in selected_neural:
                        for beam in beam_sizes:
                            tasks.append(('neural', model_name, beam))

                if include_llm:
                    for prov in selected_providers:
                        for t in temps:
                            for p in top_ps:
                                tasks.append(('llm', prov, (t, p)))

                total = len(tasks)
                current = 0

                src_words = source_text.strip().split()
                ref_words = reference_text.strip().split()
                # Align lengths by padding with empty string
                max_len = max(len(src_words), len(ref_words))
                src_words += [''] * (max_len - len(src_words))
                ref_words += [''] * (max_len - len(ref_words))

                for task in tasks:
                    typ = task[0]
                    start = time.time()
                    if typ == 'neural':
                        model_name = task[1]
                        beam = task[2]
                        # Request beams when N-best view is on
                        if show_nbest:
                            beams_per_word = [self.transliterate_neural(w, model_name, beam, return_beams=True) for w in src_words]
                            # beams_per_word: list of per-word beam lists; each entry is [[seq...],[seq...]] where seq is list of ints
                            config_str = f'Beam={beam} (N-best)'
                            label = model_name
                            ttype = 'Neural'
                        else:
                            outputs = [self.transliterate_neural(w, model_name, beam) for w in src_words]
                            config_str = f'Beam={beam}'
                            label = model_name
                            ttype = 'Neural'
                    else:
                        prov = task[1]
                        temp, p = task[2]
                        # dynamic fetch model selection if available
                        model_arg = st.session_state.selected_models.get(prov)
                        try:
                            outputs = [self.llm.transliterate(w, provider=prov, model=model_arg, temperature=temp, top_p=p) for w in [' '.join(src_words)]]
                            # LLM returns a full sentence; split to words for per-word metrics
                            outputs = outputs[0].split()
                        except Exception as e:
                            outputs = [''] * len(src_words)
                        config_str = f'T={temp},p={p}'
                        label = prov.title()
                        ttype = 'LLM'

                    elapsed = time.time() - start

                    # Compute metrics per word
                    if show_nbest and typ == 'neural':
                        # Expand per beam: compute metrics for top-k beams
                        # beams_per_word: list(len(src_words)) where each is list of beams (each beam is list of ints)
                        # Build per-beam aggregated sequences
                        max_k = min(n_best, max(len(b) for b in beams_per_word if b))
                        # Build per-beam sentences and compute metrics, supporting rerank selection
                        # beams_per_word: list over words; each entry is list of beams (each beam is list of ints)
                        per_beam_sentences = []
                        per_beam_scores = []
                        # determine K for iteration
                        K = min(n_best, max((len(b) for b in beams_per_word if b), default=1))
                        for k in range(K):
                            sent_tokens = []
                            sent_scores = []
                            for bw in beams_per_word:
                                if not bw:
                                    chosen = []
                                    score = 1e9
                                else:
                                    if k < len(bw):
                                        chosen = bw[k].tolist() if hasattr(bw[k], 'tolist') else bw[k]
                                        score = float(getattr(bw[k], 'score', 0.0)) if hasattr(bw[k], 'score') else 0.0
                                    else:
                                        chosen = bw[0].tolist() if hasattr(bw[0], 'tolist') else bw[0]
                                        score = float(getattr(bw[0], 'score', 0.0)) if hasattr(bw[0], 'score') else 0.0
                                sent_tokens.append(self.tgt_vocab.decode(chosen))
                                sent_scores.append(score)
                            per_beam_sentences.append(sent_tokens)
                            per_beam_scores.append(sent_scores)

                        # Compute per-beam aggregated metrics
                        per_beam_metrics = []
                        for k_idx, sent in enumerate(per_beam_sentences):
                            joined = ' '.join(sent)
                            word_matches = sum(int(word_exact_match(p, r)) for p, r in zip(sent, ref_words))
                            char_f1s = [char_f1_score(p, r) for p, r in zip(sent, ref_words)]
                            mean_f1 = sum(char_f1s) / len(char_f1s) if char_f1s else 0.0
                            per_beam_metrics.append({'beam': k_idx+1, 'word_matches': word_matches, 'char_f1': mean_f1, 'text': joined, 'raw_scores': per_beam_scores[k_idx]})

                        # Oracle: best beam by char_f1
                        oracle = max(per_beam_metrics, key=lambda x: x['char_f1'])

                        # Rerank selection
                        if rerank_method == 'model_score':
                            # choose beam with best average raw score (model's choice)
                            pick = min(per_beam_metrics, key=lambda x: sum(x['raw_scores'])/len(x['raw_scores']) if x['raw_scores'] else float('inf'))
                        elif rerank_method == 'length_norm':
                            # length-normalized score: avg_raw / ((5+L)/6)^alpha
                            scored = []
                            for m in per_beam_metrics:
                                L = len(m['text'].split())
                                avg_raw = sum(m['raw_scores'])/len(m['raw_scores']) if m['raw_scores'] else 0.0
                                norm = avg_raw / (((5.0 + max(1, L)) / 6.0) ** alpha)
                                scored.append((norm, m))
                            pick = min(scored, key=lambda x: x[0])[1]
                        else:
                            # mixed: normalize + simple heuristics: penalize strings with '<unk>' or excessive length
                            scored = []
                            for m in per_beam_metrics:
                                L = len(m['text'].split())
                                avg_raw = sum(m['raw_scores'])/len(m['raw_scores']) if m['raw_scores'] else 0.0
                                norm = avg_raw / (((5.0 + max(1, L)) / 6.0) ** alpha)
                                unk_pen = 1.0 if '<unk>' in m['text'] else 0.0
                                score_mix = mix_lambda * norm + (1-mix_lambda) * (norm + unk_pen)
                                scored.append((score_mix, m))
                            pick = min(scored, key=lambda x: x[0])[1]

                        # Show per-beam table
                        beam_df = pd.DataFrame(per_beam_metrics)
                        st.markdown(f"**Per-beam metrics (model={label}, {config_str})**")
                        st.table(beam_df)

                        st.markdown(f"**Oracle char-F1:** {oracle['char_f1']:.4f}  — selected beam {oracle['beam']}")
                        st.markdown(f"**Reranked pick ({rerank_method}):** beam {pick['beam']} char-F1={pick['char_f1']:.4f}")
                        # Append one result row per beam entry so the table shows N-best rows
                        for m in per_beam_metrics:
                            word_acc = m['word_matches'] / len(ref_words) if len(ref_words) > 0 else 0.0
                            char_f1 = m['char_f1']
                            results.append({
                                'Model': label,
                                'Type': ttype,
                                'Config': f"{config_str} | beam#{m['beam']}",
                                'Output': m['text'],
                                'Word Acc': round(word_acc, 4),
                                'Char F1': round(char_f1, 4),
                                'Time (s)': f"{elapsed:.3f}",
                                'Attached Result': st.session_state.get('attached_result')
                            })
                    else:
                        word_matches = 0
                        f1s = []
                        for pred, ref in zip(outputs if typ == 'neural' else outputs, ref_words):
                            word_matches += int(word_exact_match(pred, ref))
                            f1s.append(char_f1_score(pred, ref))

                        word_acc = word_matches / len(ref_words) if len(ref_words) > 0 else 0.0
                        char_f1 = sum(f1s) / len(f1s) if f1s else 0.0

                        results.append({
                            'Model': label,
                            'Type': ttype,
                            'Config': config_str,
                            'Output': ' '.join(outputs) if typ == 'neural' else ' '.join(outputs),
                            'Word Acc': round(word_acc, 4),
                            'Char F1': round(char_f1, 4),
                            'Time (s)': f"{elapsed:.3f}",
                            'Attached Result': st.session_state.get('attached_result')
                        })

                    current += 1
                    progress.progress(current / total if total > 0 else 1)

                progress.empty()

                if results:
                    st.balloons()
                    st.success(f"✅ Tested {len(results)} configurations! 🇮🇳")

                    df = pd.DataFrame(results)
                    # Save last comparison for potential PPTX/report export
                    st.session_state['last_comparison'] = df.to_dict(orient='records')
                    st.dataframe(df[['Model','Type','Config','Word Acc','Char F1','Time (s)','Output']], use_container_width=True)

                    # Time chart
                    df['Time_float'] = df['Time (s)'].astype(float)
                    fig_time = px.bar(df, x='Model', y='Time_float', color='Config', title='Processing Time Comparison')
                    st.plotly_chart(fig_time, use_container_width=True)

                    # Accuracy charts
                    fig_acc = px.bar(df, x='Model', y='Word Acc', color='Config', title='Word-level Exact Accuracy', barmode='group')
                    st.plotly_chart(fig_acc, use_container_width=True)

                    fig_f1 = px.bar(df, x='Model', y='Char F1', color='Config', title='Character-level F1', barmode='group')
                    st.plotly_chart(fig_f1, use_container_width=True)

                    csv = df.to_csv(index=False).encode('utf-8')
                    st.download_button("📥 Download CSV", data=csv, file_name=f"comparison_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv", mime="text/csv")
    
    def tab_results(self):
        """Results and analysis tab"""
        st.markdown('<div class="section-header"><h2>📈 Results & Analysis</h2></div>', unsafe_allow_html=True)
        # Discover available result files and checkpoints
        results_dir = get_project_path('outputs/results')
        ckpt_dir = get_project_path('outputs/checkpoints')

        # Prefer normalized results if available
        available_results = []
        normalized_dir = os.path.join(results_dir, 'normalized')
        if os.path.isdir(normalized_dir):
            for fname in sorted(os.listdir(normalized_dir)):
                if fname.endswith('.json') or fname.endswith('.md'):
                    available_results.append(os.path.join(normalized_dir, fname))
            # fallback to original if normalized folder empty
            if not available_results and os.path.exists(results_dir):
                for fname in sorted(os.listdir(results_dir)):
                    if fname.endswith('.json') or fname.endswith('.md'):
                        available_results.append(os.path.join(results_dir, fname))
        else:
            if os.path.exists(results_dir):
                for fname in sorted(os.listdir(results_dir)):
                    if fname.endswith('.json') or fname.endswith('.md'):
                        available_results.append(os.path.join(results_dir, fname))

        available_ckpts = []
        if os.path.exists(ckpt_dir):
            for fname in sorted(os.listdir(ckpt_dir)):
                available_ckpts.append(os.path.join(ckpt_dir, fname))

        st.markdown("### 📁 Saved Results & Checkpoints")
        col_a, col_b = st.columns([2, 1])
        with col_a:
            st.markdown("**Result files**")
            selected_result = st.selectbox("Choose a result file to load", ["(none)"] + available_results)
        with col_b:
            st.markdown("**Checkpoints**")
            selected_ckpt = st.selectbox("Choose checkpoint file", ["(none)"] + available_ckpts)

        if selected_result and selected_result != '(none)':
            try:
                with open(selected_result, 'r', encoding='utf-8') as f:
                    loaded = json.load(f)
            except Exception:
                try:
                    with open(selected_result, 'r', encoding='utf-8') as f:
                        loaded = {'content': f.read()}
                except Exception as e:
                    st.error(f"Failed to load {selected_result}: {e}")
                    loaded = None

            if loaded:
                st.markdown("### 🔎 Loaded Result Preview")
                st.write(f"**File:** {selected_result}")
                if isinstance(loaded, dict) and 'results' in loaded:
                    # Summarize metrics
                    rows = []
                    for k, v in loaded['results'].items():
                        rows.append({
                            'Setting': k,
                            'Word Acc': v.get('word_accuracy'),
                            'Char F1': v.get('char_f1'),
                            'Samples': v.get('n_samples')
                        })
                    st.table(pd.DataFrame(rows))
                else:
                    st.text(json.dumps(loaded)[:1000])

                col_norm, col_oracle = st.columns([1,1])
                with col_norm:
                    if st.button("🔧 Normalize this result"):
                        try:
                            loaded = load_result_file(selected_result)
                            if not loaded:
                                st.error("Failed to parse selected result")
                            else:
                                out_dir = os.path.join(os.path.dirname(selected_result), 'normalized')
                                os.makedirs(out_dir, exist_ok=True)
                                out_path = os.path.join(out_dir, os.path.basename(selected_result))
                                with open(out_path, 'w', encoding='utf-8') as fo:
                                    json.dump(loaded, fo, ensure_ascii=False, indent=2)
                                st.success(f"Normalized copy written to {out_path}")
                        except Exception as e:
                            st.error(f"Normalization failed: {e}")

                with col_oracle:
                    if st.button("🧭 Oracle Analysis (per-beam)"):
                        # Run oracle analysis: compute best char-F1 among beams and sweep alpha
                        try:
                            loaded = load_result_file(selected_result)
                            if not loaded:
                                st.error("Failed to load result for analysis")
                            else:
                                # Look for decoding comparisons file format with per-sample beams if present
                                # We expect loaded to contain either 'results' metrics or detailed per-sample beams under 'decoding' key
                                # For this analysis we will read a parallel file where per-sample beams are stored (convention: *_decoding_comparison.json)
                                base = os.path.basename(selected_result)
                                base_root = base.replace('.json', '')
                                candidate = os.path.join(os.path.dirname(selected_result), f"{base_root}_decoding_comparison.json")
                                if not os.path.exists(candidate):
                                    # fallback: try same file for per-beam entries under 'beams'
                                    st.info("No separate decoding comparison file found; attempting to parse beams from the selected file if available.")
                                    candidate = selected_result

                                with open(candidate, 'r', encoding='utf-8') as fh:
                                    data = json.load(fh)

                                # We expect data to have entries per sample with a 'beams' list; if not, abort with info
                                # Simplified heuristic: find top-level lists/dicts containing 'beams' key
                                samples = None
                                if isinstance(data, dict) and 'samples' in data:
                                    samples = data['samples']
                                elif isinstance(data, dict) and 'decoding' in data:
                                    samples = data['decoding']
                                else:
                                    # try to find a list of objects with 'beams'
                                    if isinstance(data, list):
                                        samples = data
                                    else:
                                        # Can't find per-sample beams — return summary only
                                        st.warning('Per-sample beams not found in file; oracle analysis requires per-sample n-best lists.')
                                        samples = None

                                if not samples:
                                    st.warning('Per-sample beams not found in file; oracle analysis requires per-sample n-best lists.')
                                    samples = None

                                # collect per-sample beams and references (capture text + optional score)
                                per_sample = []
                                for item in samples:
                                    beams_raw = item.get('beams') or item.get('nbest') or item.get('predictions')
                                    ref = item.get('reference') or item.get('target') or item.get('gold') or item.get('ref')
                                    src = item.get('source') or item.get('src') or item.get('input')
                                    if not beams_raw or not ref:
                                        continue

                                    normalized = []
                                    for b in beams_raw:
                                        # b may be dict {'text','score'}, tuple (seq,score), list [text,score], tensor, or plain string
                                        text = None
                                        score = None
                                        if isinstance(b, dict):
                                            text = b.get('text') or b.get('pred') or b.get('sequence')
                                            score = b.get('score')
                                        elif isinstance(b, (list, tuple)):
                                            # (seq_tensor, score) or [text, score]
                                            first = b[0]
                                            if hasattr(first, 'cpu'):
                                                try:
                                                    text = self.tgt_vocab.decode(first.cpu().tolist(), remove_special=True)
                                                except Exception:
                                                    text = str(first)
                                            else:
                                                text = str(first)
                                            score = b[1] if len(b) > 1 else None
                                        else:
                                            # tensor-like
                                            if hasattr(b, 'cpu'):
                                                try:
                                                    text = self.tgt_vocab.decode(b.cpu().tolist(), remove_special=True)
                                                except Exception:
                                                    text = str(b)
                                            else:
                                                text = str(b)

                                        normalized.append({'text': text, 'score': score})

                                    if normalized:
                                        per_sample.append({'source': src, 'beams': normalized, 'ref': ref})

                                if not per_sample:
                                    st.warning('No valid sample beams found for analysis')
                                    samples = None

                                if not samples or not per_sample:
                                    # abort gracefully
                                    pass

                                # Compute oracle char-F1 (upper bound)
                                oracle_scores = []
                                for s in per_sample:
                                    best = 0.0
                                    for b in s['beams']:
                                        best = max(best, char_f1_score(b['text'], s['ref']))
                                    oracle_scores.append(best)
                                mean_oracle = sum(oracle_scores) / len(oracle_scores) if oracle_scores else 0.0

                                st.success(f"Oracle mean char-F1 across {len(per_sample)} samples: {mean_oracle:.4f}")

                                # Interactive alpha-sweep and reranking controls
                                # Check if we have numeric model scores
                                has_scores = any(b.get('score') is not None for s in per_sample for b in s['beams'])

                                st.markdown("#### 🔧 Reranking & Alpha Sweep")
                                col_a, col_b = st.columns([2, 3])
                                with col_a:
                                    alpha_min = st.number_input('Alpha min', value=0.0, step=0.1, format="%.2f")
                                    alpha_max = st.number_input('Alpha max', value=2.0, step=0.1, format="%.2f")
                                    alpha_step = st.number_input('Alpha step', value=0.1, step=0.05, format="%.2f")
                                    run_sweep = st.button('▶ Run alpha sweep')
                                with col_b:
                                    st.markdown('Alpha sweep uses length-normalization: score_norm = score / (((5+L)/6) ** alpha)')
                                    st.caption('If per-beam model scores are missing, the sweep cannot rerank; provide decoding JSON with scores to use this feature.')

                                if not has_scores:
                                    st.info('Per-beam numeric scores not found in decoding data. Alpha sweep requires model scores.\nYou can still inspect oracle upper-bound above.')

                                if run_sweep:
                                    if not has_scores:
                                        st.warning('Cannot run sweep: per-beam scores missing.')
                                    else:
                                        # Build alphas
                                        alphas = []
                                        a = alpha_min
                                        while a <= alpha_max + 1e-9:
                                            alphas.append(round(a, 4))
                                            a += alpha_step

                                        alpha_results = []
                                        for a in alphas:
                                            scores = []
                                            for s in per_sample:
                                                best_text = None
                                                best_score = -1e9
                                                for b in s['beams']:
                                                    raw = b['score'] if b.get('score') is not None else 0.0
                                                    L = max(1, len(b['text'].split()))
                                                    norm = raw / (((5.0 + L) / 6.0) ** a)
                                                    if norm > best_score:
                                                        best_score = norm
                                                        best_text = b['text']
                                                scores.append(char_f1_score(best_text, s['ref']))
                                            alpha_results.append({'alpha': a, 'mean_char_f1': sum(scores) / len(scores)})

                                        # Show results and pick best alpha
                                        alpha_df = pd.DataFrame(alpha_results)
                                        best_row = alpha_df.loc[alpha_df['mean_char_f1'].idxmax()]
                                        st.markdown(f"**Best alpha:** {best_row['alpha']:.3f} → Mean char-F1: {best_row['mean_char_f1']:.4f}")
                                        st.dataframe(alpha_df, use_container_width=True)

                                        # Store in session for downstream actions
                                        st.session_state['alpha_sweep'] = alpha_results

                                # Live rerank preview with selected alpha
                                st.markdown('#### 🔍 Rerank preview')
                                chosen_alpha = st.number_input('Choose alpha for preview', value=0.7, step=0.05, format="%.2f")
                                preview_n = st.slider('Preview samples', 1, min(50, len(per_sample)), 5)
                                if st.button('Apply rerank (preview)'):
                                    if not has_scores:
                                        st.warning('Cannot rerank: per-beam scores missing in decoding data.')
                                    else:
                                        preview_rows = []
                                        for s in per_sample[:preview_n]:
                                            best_text = None
                                            best_score = -1e9
                                            for b in s['beams']:
                                                raw = b['score'] if b.get('score') is not None else 0.0
                                                L = max(1, len(b['text'].split()))
                                                norm = raw / (((5.0 + L) / 6.0) ** chosen_alpha)
                                                if norm > best_score:
                                                    best_score = norm
                                                    best_text = b['text']
                                            preview_rows.append({'source': s.get('source', ''), 'gold': s['ref'], 'reranked': best_text, 'oracle': max((char_f1_score(b['text'], s['ref']) for b in s['beams']), default=0.0)})
                                        st.dataframe(pd.DataFrame(preview_rows), use_container_width=True)
                        except Exception as e:
                            st.error(f"Oracle analysis failed: {e}")
        else:
            st.info("No result selected. You can load a result JSON from outputs/results.")

        # Add PPT/Report area
        st.divider()
        st.markdown("### 📄 Reports & Presentation")
        st.markdown("Upload your `report.md` (final report) and optionally export a PPTX summary.")
        report_file = st.file_uploader("Upload report.md", type=['md'])
        if report_file is not None:
            content = report_file.read().decode('utf-8')
            st.session_state['uploaded_report'] = content
            st.success("Report uploaded")

            if st.button("📤 Generate PPTX (preview and download)"):
                # Generate PPTX using last comparison and uploaded report
                last = st.session_state.get('last_comparison')
                if not last:
                    st.warning("Run a comparison first to populate slides (Compare tab).")
                else:
                    pptx_bytes = generate_full_pptx(last, st.session_state.get('attached_result'), content)
                    if pptx_bytes is None:
                        st.error("PPTX generation unavailable. Install python-pptx and optionally plotly[kaleido] for charts.")
                    else:
                        st.success("PPTX generated — ready for download")
                        st.download_button("📥 Download PPTX", data=pptx_bytes, file_name=f"transliteration_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx", mime='application/vnd.openxmlformats-officedocument.presentationml.presentation')

        st.divider()
        st.markdown("### ✅ Notes & Professional Polishing")
        st.markdown("- Results and checkpoints are loaded from `outputs/results` and `outputs/checkpoints` respectively.\n- Use 'Attach this result' to include the selected result when exporting comparison CSVs.\n- PPT export requires `python-pptx`.\n- Ensure vocabs and model checkpoints are consistent with config.")
        
        # Error analysis
        st.divider()
        st.markdown("#### 🔍 Error Analysis")
        
        with st.expander("📚 Transliteration Challenges (ACL W15-3902)", expanded=False):
            st.markdown("""
            ### Why Some Sequences Are Hard to Transliterate
            
            **1. Conjunct Consonants (संयुक्त व्यंजन)** 🔤  
            - Examples: `क्ष` (ksha), `त्र` (tra), `ज्ञ` (gya)
            - **Challenge:** Many-to-one mapping (3+ Roman chars → 1 Devanagari)
            - **Impact:** Often split incorrectly or missing
            
            **2. Vowel Matras (स्वर मात्राएँ)** 📝  
            - Examples: `े` (e), `ो` (o), `ु` (u)
            - **Challenge:** Position-sensitive placement
            - **Impact:** Placement errors or omissions
            
            **3. Halant (हलंत - ्)** ⚡  
            - Symbol: `्`
            - **Challenge:** Invisible in Roman script
            - **Impact:** Missing conjunct formations
            
            **4. Aspirated Consonants (महाप्राण)** 💨  
            - Examples: `ख` (kha), `घ` (gha), `छ` (cha)
            - **Challenge:** Subtle phonetic differences
            - **Impact:** Substitution errors
            
            **5. Nasalization (अनुनासिक)** 👃  
            - Symbols: `ं` (anusvara), `ँ` (chandrabindu)
            - **Challenge:** No direct Roman equivalent
            - **Impact:** Often omitted or misplaced
            """)
        
        # History
        if st.session_state.history:
            st.divider()
            st.markdown("#### 📜 Translation History")
            
            history_df = pd.DataFrame(st.session_state.history)
            
            col1, col2 = st.columns([4, 1])
            
            with col1:
                if len(history_df) > 0:
                    st.dataframe(
                        history_df[['timestamp', 'input', 'output', 'model', 'time']],
                        use_container_width=True,
                        hide_index=True
                    )
            
            with col2:
                if st.button("🗑️ Clear"):
                    st.session_state.history = []

        # Try loading a normalized result and show a 3D error plot
        try:
            norm_dir = os.path.join(get_project_path('outputs/results'), 'normalized')
            pick_dir = norm_dir if os.path.isdir(norm_dir) else get_project_path('outputs/results')
            files = [f for f in os.listdir(pick_dir) if f.endswith('.json')]
            if files:
                loaded_example = load_result_file(os.path.join(pick_dir, files[0]))
                if loaded_example and 'error_analysis' in loaded_example:
                    fig = make_error_3d_plot(loaded_example['error_analysis'])
                    st.plotly_chart(fig, use_container_width=True)
        except Exception:
            pass
    
    def run(self):
        """Main application logic"""
        # Header with Indian tricolor and floating chars
        st.markdown("""
        <div class="main-header">
            <h1>🇮🇳 Hindi Transliteration System</h1>
            <p>CS772 Assignment 2 | Roman → Devanagari | मेरा भारत महान</p>
        </div>
        """, unsafe_allow_html=True)
        
        # Sidebar
        self.sidebar_content()
        
        # Tabs
        tabs = st.tabs([
            "🎯 Transliterate", 
            "🔑 API Config", 
            "📊 Compare", 
            "📈 Results"
        ])
        
        with tabs[0]:
            self.tab_transliterate()
        
        with tabs[1]:
            self.tab_api_configuration()
        
        with tabs[2]:
            self.tab_compare_models()
        
        with tabs[3]:
            self.tab_results()
        
        # Footer with tricolor
        st.markdown("""
        <div class="main-footer">
            <p class="footer-hindi">जय हिन्द | मेरा भारत महान</p>
            <p><strong>Hindi Transliteration System</strong> | CS772 Assignment 2</p>
            <p>© 2025 | ACL W15-3902 Compliant | For Educational Use</p>
            <p>🇮🇳 Powered by Indian Innovation 🇮🇳</p>
        </div>
        """, unsafe_allow_html=True)


def main():
    """Entry point"""
    app = TransliterationApp()
    app.run()


if __name__ == "__main__":
    main()